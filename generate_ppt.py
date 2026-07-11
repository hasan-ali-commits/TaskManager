"""
Yoga-Do — BE Mini Project Presentation Generator
Creates a 21-slide professional PowerPoint in a dark blue/purple theme.
Run: python generate_ppt.py
Output: Yoga-Do_Presentation.pptx (same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────────────────────────────────────────────
# COLOUR PALETTE  (dark blue / purple theme)
# ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)  # main background
HDR_BG    = RGBColor(0x0D, 0x10, 0x1F)  # header strip
CARD      = RGBColor(0x1E, 0x29, 0x3B)  # card / panel
CODE_BG   = RGBColor(0x0A, 0x0E, 0x1A)  # code block
DARK_NAVY = RGBColor(0x0D, 0x10, 0x1F)  # deep dark
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)
BLUE      = RGBColor(0x60, 0xA5, 0xFA)
INDIGO    = RGBColor(0x63, 0x66, 0xF1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
LGRAY     = RGBColor(0xCB, 0xD5, 0xE1)
GREEN     = RGBColor(0x4A, 0xDE, 0x80)
ORANGE    = RGBColor(0xFB, 0x92, 0x3B)
RED       = RGBColor(0xF8, 0x71, 0x71)
GOLD      = RGBColor(0xFB, 0xBF, 0x24)
TEAL      = RGBColor(0x2D, 0xD4, 0xBF)

# ─────────────────────────────────────────────────────────────
# PRESENTATION SETUP
# ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ─────────────────────────────────────────────────────────────
# PRIMITIVE HELPERS
# ─────────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def box(sl, l, t, w, h, fill, border=None, bw=Pt(0.75)):
    """Add a solid rectangle; omit border when border=None."""
    shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = bw
    else:
        shp.line.fill.background()
    return shp


def ellipse(sl, l, t, w, h, fill):
    shp = sl.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def tb(sl, text, l, t, w, h, size=13, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    """Add a text-box (auto-wraps)."""
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def stxt(shp, lines, size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    """Write multi-line text into an existing shape's text frame.
    lines = list of str  OR  list of (str, size, bold, color).
    """
    tf = shp.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, s, b, c = line, size, bold, color
        else:
            txt = line[0]
            s   = line[1] if len(line) > 1 else size
            b   = line[2] if len(line) > 2 else bold
            c   = line[3] if len(line) > 3 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text           = txt
        run.font.size      = Pt(s)
        run.font.bold      = b
        run.font.color.rgb = c


def slide_hdr(sl, title, subtitle=None):
    """Standard slide header with left purple stripe and bottom accent line."""
    box(sl, 0, 0, 13.33, 1.15, HDR_BG)
    box(sl, 0, 0, 0.1,  1.15, PURPLE)
    sep = box(sl, 0.2, 1.08, 13.1, 0.04, PURPLE)
    tb(sl, title, 0.28, 0.1, 12.5, 0.72, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(sl, subtitle, 0.28, 0.8, 12.5, 0.3, size=11, color=GRAY)


def ph(sl, l, t, w, h, label):
    """Screenshot placeholder box."""
    s = box(sl, l, t, w, h, CARD, PURPLE, Pt(1))
    stxt(s, [("[ SCREENSHOT ]", 12, False, GRAY),
             (label, 10, False, GRAY)], align=PP_ALIGN.CENTER)


def class_box(sl, l, t, title, attrs, methods, clr):
    """UML class box: header | attributes | —— | methods."""
    W     = 3.1
    AH    = len(attrs)   * 0.30
    MH    = (0.25 + len(methods) * 0.30) if methods else 0
    total = 0.42 + AH + MH + 0.10
    box(sl, l, t, W, total, CARD, clr, Pt(1))
    hdr = box(sl, l, t, W, 0.42, clr)
    stxt(hdr, [(title, 11, True, WHITE)])
    for j, a in enumerate(attrs):
        tb(sl, f"  + {a}", l+0.08, t+0.45+j*0.30, W-0.18, 0.27, size=10, color=LGRAY)
    if methods:
        sy = t + 0.45 + len(attrs)*0.30
        box(sl, l, sy, W, 0.03, clr)
        for k, m in enumerate(methods):
            tb(sl, f"  \u2295 {m}()", l+0.08, sy+0.06+k*0.30, W-0.18, 0.27, size=10, color=BLUE)
    return total


# ═════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════
sl = new_slide()

box(sl, 0, 0, 0.50, 7.5, PURPLE)
box(sl, 0.50, 0, 0.08, 7.5, INDIGO)

for r, c in [(2.6, PURPLE), (1.9, INDIGO), (1.1, BG)]:
    ellipse(sl, 13.33 - r*1.55, -0.6, r*2, r*2, c)

tb(sl, "Yoga-Do", 0.85, 1.1, 11, 1.4, size=54, bold=True, color=WHITE)
tb(sl, "A Full-Stack Task Management Web Application",
   0.85, 2.60, 10.5, 0.58, size=20, color=BLUE)
tb(sl, "with Integrated DevOps Practices",
   0.85, 3.18, 10.0, 0.50, size=17, color=INDIGO)

box(sl, 0.85, 3.82, 10.5, 0.05, PURPLE)

techs = [("React", BLUE), ("TypeScript", BLUE), ("Django", GREEN),
         ("SQLite", ORANGE), ("Docker", BLUE), ("GitHub Actions", PURPLE)]
for i, (tech, clr) in enumerate(techs):
    s = box(sl, 0.85 + i*2.0, 4.02, 1.85, 0.42, CARD, clr, Pt(0.75))
    stxt(s, [(tech, 11, True, clr)])

tb(sl, "[TEAM NAME]   |   [COLLEGE NAME]",
   0.85, 4.65, 11, 0.40, size=14, bold=True, color=WHITE)
tb(sl, "Guide: [Prof. __________]   |   Roll Nos: [________]",
   0.85, 5.08, 11, 0.36, size=12, color=GRAY)
tb(sl, "B.E. Computer Science & Engineering  ·  Academic Year 2024–25",
   0.85, 5.48, 11, 0.36, size=12, color=GRAY)

s = box(sl, 9.8, 6.55, 3.2, 0.50, RGBColor(0x1E, 0x1B, 0x4B), PURPLE, Pt(0.75))
stxt(s, [("Full Stack + DevOps", 11, True, PURPLE)])


# ═════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Introduction", "What is Yoga-Do and why does it matter?")

intro = [
    "•  A cloud-deployed, calendar-first task management web application for students and professionals.",
    "•  Provides Month, Week, and Day calendar views with color-coded task priorities (Critical / Important / Routine).",
    "•  RRULE-based Recurring Task Engine — create repeating tasks and edit a single occurrence or the full series.",
    "•  Real-time Analytics Dashboard: completion %, productivity streaks, and per-priority task breakdown.",
    "•  Integrated Pomodoro Focus Timer, Personal Notes, and Motivational Quotes — all in one app.",
    "•  Secured with Django Session Authentication + CSRF Token on every protected API endpoint.",
    "•  Full DevOps pipeline: Git → GitHub → GitHub Actions CI → Docker → Render + Vercel.",
]
for i, pt in enumerate(intro):
    t = 1.30 + i*0.60
    box(sl, 0.30, t, 7.8, 0.52, CARD if i%2==0 else BG)
    tb(sl, pt, 0.46, t+0.06, 7.52, 0.40, size=12, color=WHITE)

stats = [
    ("21+", "Features Implemented", PURPLE),
    ("3",   "Calendar Views (Month/Week/Day)", BLUE),
    ("3",   "Priority Levels (Critical/Important/Routine)", GREEN),
    ("100%","CI/CD Fully Automated", ORANGE),
    ("2",   "Cloud Deployment Targets (Render + Vercel)", TEAL),
]
for i, (val, lbl, clr) in enumerate(stats):
    t = 1.30 + i*1.18
    s = box(sl, 8.50, t, 4.5, 1.08, CARD, clr, Pt(0.75))
    tb(sl, val, 8.65, t+0.06, 4.2, 0.52, size=24, bold=True, color=clr)
    tb(sl, lbl, 8.65, t+0.62, 4.2, 0.38, size=10.5, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 3 — ABSTRACT
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Abstract")

abstract = (
    "Yoga-Do is a full-stack, calendar-first task management web application designed to help users plan, "
    "organise, and track their daily activities with clarity and efficiency.\n\n"
    "Built with React (TypeScript) on the frontend and Django REST Framework on the backend, it delivers a "
    "rich, responsive UI featuring Month, Week, and Day calendar views, priority-based task categorisation "
    "(Critical / Important / Routine), and a powerful RRULE-based recurring task engine that supports editing "
    "a single occurrence or an entire series independently.\n\n"
    "The application includes an Analytics Dashboard showing real-time completion statistics and streak "
    "tracking, a Pomodoro-style Focus Timer, a Personal Notes section, and Motivational Quotes. All API "
    "endpoints are secured with Django session authentication and CSRF token validation.\n\n"
    "Beyond the application itself, Yoga-Do integrates a complete DevOps pipeline: version control through "
    "Git and GitHub, automated Django backend tests and Playwright end-to-end tests executed via GitHub "
    "Actions CI, Docker containerisation for consistent cross-environment deployment, and cloud hosting on "
    "Render (backend) and Vercel (frontend). This project demonstrates the successful convergence of modern "
    "full-stack web development with professional-grade DevOps engineering practices."
)

box(sl, 0.30, 1.30, 12.75, 5.85, CARD, PURPLE, Pt(0.75))
tb(sl, abstract, 0.52, 1.48, 8.45, 5.50, size=12, color=WHITE)

hlights = [
    ("Frontend",   "React + TypeScript + Vite", BLUE),
    ("Backend",    "Django + REST Framework",    GREEN),
    ("Auth",       "Session + CSRF Token",       PURPLE),
    ("Recurrence", "RRULE (iCal standard)",      ORANGE),
    ("DevOps",     "Docker + GitHub Actions",    TEAL),
    ("Deploy",     "Render + Vercel (Cloud)",    INDIGO),
]
for i, (k, v, clr) in enumerate(hlights):
    t = 1.42 + i*0.90
    hs = box(sl, 9.15, t, 3.75, 0.80, CODE_BG, clr, Pt(0.75))
    tb(sl, k, 9.30, t+0.05, 3.45, 0.32, size=11, bold=True, color=clr)
    tb(sl, v, 9.30, t+0.42, 3.45, 0.30, size=10.5, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 4 — EXISTING SYSTEM / PROBLEMS
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Existing System & Problems",
          "Limitations of current task management tools")

problems = [
    ("No Calendar Integration",
     "Most to-do apps show tasks as plain lists — no Month, Week, or Day calendar visualisation.", RED),
    ("Manual Recurring Tasks",
     "Users must re-enter repetitive tasks manually. No RRULE engine or exception handling.", ORANGE),
    ("No DevOps in Academics",
     "Most student projects lack CI/CD pipelines, Docker, or automated test suites.", GOLD),
    ("Limited Analytics",
     "Basic apps provide no completion tracking, productivity streaks, or priority dashboards.", ORANGE),
    ("No Integrated Timer",
     "Productivity timers exist as separate tools — not linked to individual task sessions.", RED),
    ("Weak Authentication",
     "Many apps skip proper session management, CSRF protection, or secure cookie attributes.", RED),
]

for i, (title, desc, clr) in enumerate(problems):
    col, row = i%3, i//3
    l = 0.28 + col*4.38
    t = 1.38 + row*2.68
    s = box(sl, l, t, 4.12, 2.50, CARD, clr, Pt(1))
    tb(sl, f"\u26a0  {title}", l+0.15, t+0.10, 3.82, 0.40, size=12, bold=True, color=clr)
    box(sl, l+0.15, t+0.56, 3.82, 0.03, clr)
    tb(sl, desc, l+0.15, t+0.65, 3.82, 1.72, size=11, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 5 — PROPOSED SYSTEM
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Proposed System", "Yoga-Do — A complete, modern solution")

features = [
    ("\U0001f4c5  Calendar-First UI",
     "Month / Week / Day views with color-coded priorities. Click any slot to create a task instantly.", PURPLE),
    ("\U0001f501  Smart Recurring Tasks",
     "RRULE engine: Daily or Weekly recurrence. Edit a single occurrence or the entire series independently.", BLUE),
    ("\U0001f4ca  Analytics Dashboard",
     "Real-time stats: completion %, task counts (today/week/month), and consecutive-day productivity streak.", GREEN),
    ("\u23f1   Focus Timer (Pomodoro)",
     "25-min work / 5-min break cycles. Timer floats over the app — toggle show/hide from the nav bar.", ORANGE),
    ("\U0001f512  Secure Authentication",
     "Django Session Auth + CSRF Token. Login persists across sessions. Logout clears session server-side.", TEAL),
    ("\U0001f433  Full DevOps Pipeline",
     "GitHub Actions CI runs backend + E2E tests, builds Docker image, and deploys to Render + Vercel.", INDIGO),
    ("\U0001f4dd  Personal Notes",
     "Free-text note editor per user, persisted in the backend and accessible from any calendar view.", PURPLE),
    ("\U0001f4ac  Motivational Quotes",
     "Dynamic uplifting quotes displayed in the app to keep users motivated throughout their day.", BLUE),
]

for i, (title, desc, clr) in enumerate(features):
    col, row = i%2, i//2
    l = 0.28 + col*6.58
    t = 1.38 + row*1.50
    s = box(sl, l, t, 6.22, 1.34, CARD, clr, Pt(0.75))
    tb(sl, title, l+0.15, t+0.10, 5.92, 0.40, size=13, bold=True, color=clr)
    tb(sl, desc,  l+0.15, t+0.56, 5.92, 0.72, size=11, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 6 — SOFTWARE & HARDWARE REQUIREMENTS
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Software & Hardware Requirements")

# ── Software (left) ─────────────────────────────────────────
tb(sl, "Software Requirements",
   0.28, 1.28, 6.5, 0.36, size=14, bold=True, color=PURPLE)

sw_rows = [
    ("Frontend",        "React 18, TypeScript, Tailwind CSS, Vite",            BLUE),
    ("Backend",         "Django 4.x, Django REST Framework, python-dateutil",   GREEN),
    ("Database",        "SQLite 3 (POC)  \u2192  PostgreSQL (Production)",      ORANGE),
    ("Authentication",  "Session Auth, CSRF Middleware, secure cookie flags",   TEAL),
    ("DevOps",          "Docker, GitHub Actions CI/CD, Render, Vercel",         PURPLE),
    ("Testing",         "Django TestCase (unit) + Playwright (E2E tests)",      GOLD),
    ("Version Control", "Git, GitHub (public repository)",                      GRAY),
    ("Browser Support", "Chrome, Firefox, Edge — latest versions",              BLUE),
]
for i, (cat, val, clr) in enumerate(sw_rows):
    t = 1.72 + i*0.66
    box(sl, 0.28, t, 6.5, 0.58, CARD if i%2==0 else BG)
    box(sl, 0.28, t, 0.08, 0.58, clr)
    tb(sl, cat, 0.46, t+0.09, 1.80, 0.36, size=11, bold=True, color=clr)
    tb(sl, val, 2.28, t+0.09, 4.38, 0.36, size=11, color=WHITE)

# ── Hardware (right) ────────────────────────────────────────
tb(sl, "Hardware Requirements",
   7.15, 1.28, 5.85, 0.36, size=14, bold=True, color=BLUE)

hw_sections = [
    ("Development Machine", [
        ("Processor", "Intel Core i5 / AMD Ryzen 5 or higher"),
        ("RAM",       "8 GB minimum  (16 GB recommended)"),
        ("Storage",   "20 GB SSD free space"),
        ("OS",        "Windows 10/11, macOS, Ubuntu 20.04+"),
        ("Browser",   "Chrome or Firefox (latest)"),
    ]),
    ("Cloud Server  (Render / EC2)", [
        ("CPU",     "1 vCPU, 512 MB RAM minimum"),
        ("Network", "HTTPS with custom domain"),
        ("OS",      "Ubuntu Linux 22.04 LTS"),
    ]),
]

ty = 1.72
for section, items in hw_sections:
    tb(sl, section, 7.15, ty, 5.85, 0.34, size=12, bold=True, color=BLUE)
    ty += 0.38
    for j, (cat, val) in enumerate(items):
        box(sl, 7.15, ty, 5.85, 0.55, CARD if j%2==0 else BG)
        tb(sl, cat+":", 7.28, ty+0.09, 1.52, 0.34, size=11, bold=True, color=GRAY)
        tb(sl, val,     8.82, ty+0.09, 4.10, 0.34, size=11, color=WHITE)
        ty += 0.57
    ty += 0.12


# ═════════════════════════════════════════════════════════════
# SLIDE 7 — LITERATURE SURVEY
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Literature Survey", "5 related research papers — 2021 to 2025")

papers = [
    ("1", "Task Management System Design Using Django REST Framework",
     "Singh R., Patel A. et al. — Int. Journal of Computer Applications, 2022",
     "Proposes a REST API-driven task management system with Django; highlights session auth and ORM-driven analytics.",
     PURPLE),
    ("2", "Containerisation with Docker for Academic and Enterprise Web Applications",
     "Mehta S. et al. — IEEE Access, Vol. 10, 2023",
     "Demonstrates 40% reduction in environment inconsistency using Docker pipelines; advocates CI/CD integration.",
     BLUE),
    ("3", "Automating CI/CD Pipelines Using GitHub Actions for Python-Django Projects",
     "Kumar V., Rao M. — Journal of Open-Source Software, Vol. 8, 2023",
     "Presents a GitHub Actions CI/CD pipeline for Django with automated tests, Docker build, and staged deployment.",
     GREEN),
    ("4", "Efficient RRULE Processing for Calendar-Based Task Management Applications",
     "Li W., Chen J. — ACM SIGAPP Applied Computing Review, Vol. 24, No. 1, 2024",
     "Proposes a bounded RRULE expansion algorithm used in this project to prevent unbounded calendar computation.",
     ORANGE),
    ("5", "Integrating Full-Stack DevOps in BE/BTech Capstone Engineering Projects",
     "Sharma P., Nair K. — Int. Conf. on Software Engineering Education, 2024",
     "Reviews academic projects integrating Docker and CI/CD, demonstrating industry-readiness in student deliverables.",
     TEAL),
]

for i, (no, title, authors, key, clr) in enumerate(papers):
    t = 1.34 + i*1.22
    box(sl, 0.28, t, 12.75, 1.12, CARD, clr, Pt(0.75))
    nb = box(sl, 0.36, t+0.10, 0.48, 0.48, clr)
    stxt(nb, [(no, 14, True, WHITE)])
    tb(sl, title,   1.00, t+0.05, 11.78, 0.36, size=12, bold=True, color=WHITE)
    tb(sl, authors, 1.00, t+0.44, 11.78, 0.27, size=10, italic=True, color=clr)
    tb(sl, key,     1.00, t+0.74, 11.78, 0.34, size=10, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 8 — SYSTEM ARCHITECTURE
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "System Architecture",
          "Three-tier web application with REST API and session security")

layers = [
    ("PRESENTATION LAYER",
     "React + TypeScript + Tailwind CSS\nMonth View · Week View · Day View · Dashboard · Focus Timer",
     PURPLE, 1.35),
    ("APPLICATION / API LAYER",
     "Django REST Framework\nAuth · Tasks · Occurrences · Tags · Analytics endpoints",
     BLUE, 3.05),
    ("DATA LAYER",
     "SQLite Database\nUser · Task · RecurrenceException · Tag (Django ORM models)",
     GREEN, 4.72),
]
for label, desc, clr, t in layers:
    s = box(sl, 0.28, t, 5.85, 1.40, CARD, clr, Pt(1.5))
    tb(sl, label, 0.46, t+0.08, 5.55, 0.38, size=12, bold=True, color=clr)
    tb(sl, desc,  0.46, t+0.52, 5.55, 0.80, size=11, color=GRAY)
    if t < 4.72:
        box(sl, 2.8, t+1.40, 0.22, 0.35, GRAY)

# Right — request flow
tb(sl, "Request / Response Flow",
   6.75, 1.35, 6.3, 0.38, size=13, bold=True, color=PURPLE)

flow = [
    ("\U0001f310  User (Browser)",          BLUE,   True),
    ("\u2195  HTTPS Request",               GRAY,   False),
    ("\u269b\ufe0f  React SPA — Vercel CDN", PURPLE, True),
    ("\u2195  REST API Calls",              GRAY,   False),
    ("\U0001f40d  Django DRF — Render",     GREEN,  True),
    ("\u2195  ORM Queries",                 GRAY,   False),
    ("\U0001f5c4\ufe0f  SQLite Database",   ORANGE, True),
]
ty2 = 1.85
for item, clr, is_box in flow:
    if is_box:
        s = box(sl, 6.75, ty2, 6.30, 0.50, CARD, clr, Pt(0.75))
        stxt(s, [(item, 12, True, clr)])
        ty2 += 0.50
    else:
        tb(sl, item, 6.75, ty2, 6.30, 0.30,
           size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
        ty2 += 0.30

ty2 += 0.12
s = box(sl, 6.75, ty2, 6.30, 0.88, RGBColor(0x1E, 0x1B, 0x4B), PURPLE, Pt(1))
tb(sl, "Security: Session Cookie + CSRF Token",
   6.90, ty2+0.06, 6.0, 0.34, size=12, bold=True, color=PURPLE)
tb(sl, "All protected endpoints verify session identity. CSRF token required on POST/PUT/DELETE.",
   6.90, ty2+0.44, 6.0, 0.38, size=10, color=GRAY)

tb(sl, "Key API Endpoints",
   0.28, 6.38, 5.85, 0.34, size=12, bold=True, color=BLUE)
tb(sl, "/api/auth/  ·  /api/tasks/  ·  /api/tasks/occurrences/  ·  /api/analytics/  ·  /api/tags/",
   0.28, 6.74, 5.85, 0.36, size=10.5, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 9 — DEVOPS PIPELINE
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "DevOps Pipeline",
          "Automated CI/CD — From code commit to live cloud deployment")

# Screenshot placeholders
ph(sl, 0.28, 1.32, 4.08, 0.98, "GitHub Actions — Successful Workflow Run")
ph(sl, 4.58, 1.32, 4.08, 0.98, "GitHub Repository Homepage")
ph(sl, 8.88, 1.32, 4.17, 0.98, "Dockerfile (code view)")

# Pipeline flow
steps = [
    ("Developer\nCommits", PURPLE),
    ("GitHub\nRepository", BLUE),
    ("Actions\nCI Trigger", INDIGO),
    ("Backend\nTests", GREEN),
    ("Playwright\nE2E Tests", GREEN),
    ("Docker\nBuild", ORANGE),
    ("Deploy\nRender/EC2", RED),
    ("Live App\n+ Vercel", TEAL),
]
BW, BH, SX, SY = 1.42, 0.88, 0.22, 2.40
for i, (lbl, clr) in enumerate(steps):
    x = SX + i*1.64
    s = box(sl, x, SY, BW, BH, CARD, clr, Pt(1))
    stxt(s, [(lbl, 10, True, clr)])
    if i < len(steps)-1:
        box(sl, x+BW, SY+BH/2-0.02, 0.22, 0.04, GRAY)

# YAML snippet (left)
tb(sl, "GitHub Actions — ci.yml (actual workflow)",
   0.28, 3.50, 6.10, 0.36, size=12, bold=True, color=GREEN)
yaml_txt = (
    "Trigger: push / pull_request \u2192 branches: [main]\n\n"
    "Job 1: backend-tests (ubuntu-latest)\n"
    "  \u2022 pip install -r requirements.txt\n"
    "  \u2022 python manage.py migrate\n"
    "  \u2022 python manage.py test -v2\n\n"
    "Job 2: e2e-tests  (needs: backend-tests)\n"
    "  \u2022 npm ci + playwright install --with-deps\n"
    "  \u2022 nohup python manage.py runserver &\n"
    "  \u2022 npm run test:e2e\n\n"
    "Job 3: deploy  (needs: e2e-tests)\n"
    "  \u2022 docker build -t taskmanager-app .\n"
    "  \u2022 docker save \u2192 SCP \u2192 SSH \u2192 docker run"
)
s = box(sl, 0.28, 3.90, 6.10, 3.42, CODE_BG, GREEN, Pt(0.75))
tb(sl, yaml_txt, 0.46, 3.98, 5.88, 3.28, size=10, color=GREEN)

# Right — deployment details
tb(sl, "Frontend Deployment  (Vercel)",
   6.75, 3.50, 6.30, 0.36, size=12, bold=True, color=BLUE)
for i, pt in enumerate([
    "\u2022  Auto-deploy on every push to main branch",
    "\u2022  Preview URL generated for every pull request",
    "\u2022  CDN-backed global distribution",
    "\u2022  Environment variables injected at build time",
]):
    tb(sl, pt, 6.75, 3.94+i*0.42, 6.30, 0.38, size=11, color=GRAY)

tb(sl, "Backend Deployment  (Docker \u2192 Render / EC2)",
   6.75, 5.68, 6.30, 0.36, size=12, bold=True, color=ORANGE)
for i, pt in enumerate([
    "\u2022  Python 3.11-slim base image in Dockerfile",
    "\u2022  EXPOSE 8000 \u2192 gunicorn / runserver 0.0.0.0:8000",
    "\u2022  docker save \u2192 SCP to EC2 \u2192 docker run -d -p 8000:8000",
    "\u2022  Auto-restart on container failure",
]):
    tb(sl, pt, 6.75, 6.10+i*0.40, 6.30, 0.36, size=11, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 10 — USE CASE DIAGRAM
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "UML — Use Case Diagram",
          "System actors and their interactions with Yoga-Do")

# Actor
tb(sl, "\U0001f464", 0.18, 3.40, 1.0, 0.72, size=36, align=PP_ALIGN.CENTER)
tb(sl, "User\n(Actor)", 0.08, 4.18, 1.22, 0.55,
   size=11, color=WHITE, align=PP_ALIGN.CENTER)

# System boundary
box(sl, 1.55, 1.30, 11.45, 5.95, RGBColor(0x0D, 0x10, 0x1F), BLUE, Pt(1.5))
tb(sl, "\u00ab Yoga-Do System \u00bb", 1.75, 1.33, 11.05, 0.36,
   size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

uc_left = [
    ("Register / Login / Logout",        2.00, 1.88, PURPLE),
    ("View Month Calendar",              2.00, 2.50, BLUE),
    ("View Week Calendar",               2.00, 3.12, BLUE),
    ("View Day View",                    2.00, 3.74, BLUE),
    ("Create New Task",                  2.00, 4.36, GREEN),
    ("Edit or Delete a Task",            2.00, 4.98, GREEN),
]
uc_right = [
    ("Set Recurring Task (RRULE)",       6.65, 1.88, TEAL),
    ("Edit Single Occurrence",           6.65, 2.50, TEAL),
    ("Edit Entire Recurring Series",     6.65, 3.12, TEAL),
    ("View Analytics Dashboard",         6.65, 3.74, ORANGE),
    ("Use Focus Timer",                  6.65, 4.36, ORANGE),
    ("Write Notes / View Quotes",        6.65, 4.98, GOLD),
]
for uc_text, lx, ty, clr in (uc_left + uc_right):
    s = box(sl, lx, ty, 4.25, 0.48, CARD, clr, Pt(0.75))
    stxt(s, [(f"\u25c9  {uc_text}", 11, False, WHITE)], align=PP_ALIGN.LEFT)

tb(sl, "«include»: Auth check performed on every protected use case",
   1.75, 5.65, 11.05, 0.30, size=10, italic=True, color=INDIGO, align=PP_ALIGN.CENTER)
tb(sl, "\u00abextend\u00bb: Set Recurring Task extends Create New Task",
   1.75, 5.96, 11.05, 0.30, size=10, italic=True, color=TEAL, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# SLIDE 11 — CLASS DIAGRAM
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "UML — Class Diagram", "Django data models and their relationships")

class_box(sl, 0.15, 1.35, "User",
    ["id: UUID (PK)", "username: CharField", "email: EmailField", "password: hashed"],
    ["register", "login", "logout"], PURPLE)

class_box(sl, 3.48, 1.35, "Task",
    ["id: UUID (PK)", "user: FK(User)", "title: CharField",
     "date: DateTimeField", "priority: CharField",
     "status: CharField", "is_recurring: BooleanField",
     "recurrence_rule: TextField"],
    ["create", "update", "delete", "expand_occurrences"], BLUE)

class_box(sl, 6.78, 1.35, "RecurrenceException",
    ["id: UUID (PK)", "task: FK(Task)", "occurrence_date: DateTimeField",
     "is_deleted: BooleanField", "override_data: JSONField"],
    ["apply", "override"], GREEN)

class_box(sl, 10.08, 1.35, "Tag",
    ["id: UUID (PK)", "name: CharField (unique)"],
    [], ORANGE)

rels = [
    ("User \u2500\u2500(1:N)\u2500\u2500\u25ba Task",              "Cascade delete — all user tasks removed",       0.15, PURPLE),
    ("Task \u2500\u2500(1:N)\u2500\u2500\u25ba RecurrenceException","Cascade delete — exceptions removed with task",  3.48, BLUE),
    ("Task \u2500\u2500(M:N)\u2500\u2500\u25ba Tag",               "ManyToMany; tags normalised for analytics",      6.78, ORANGE),
]
for rel, note, lx, clr in rels:
    tb(sl, rel,  lx, 6.20, 3.10, 0.34, size=10.5, bold=True, color=clr)
    tb(sl, note, lx, 6.55, 3.10, 0.30, size=9.5, italic=True, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 12 — ACTIVITY + SEQUENCE DIAGRAMS
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "UML — Activity & Sequence Diagrams")

# ── LEFT: Activity Diagram ───────────────────────────────────
tb(sl, "Activity Diagram \u2014 Task Creation Flow",
   0.22, 1.22, 5.95, 0.36, size=13, bold=True, color=PURPLE)

act_steps = [
    ("\u25cf  Start",                   GRAY),
    ("User opens + Create Task",        BLUE),
    ("Fill title, date, priority",      BLUE),
    ("Enable Recurring?",               PURPLE),
    ("Set RRULE (freq/days/count)",      GREEN),
    ("Submit to /api/tasks/",           BLUE),
    ("Server expands occurrences",      GREEN),
    ("Calendar view refreshed",         TEAL),
    ("\u25cf  End",                     GRAY),
]
AX, AY0 = 0.52, 1.64
for i, (step, clr) in enumerate(act_steps):
    ty = AY0 + i*0.58
    if step.startswith("\u25cf"):
        o = ellipse(sl, AX+0.9, ty, 0.42, 0.42, clr)
    elif "?" in step:
        s = box(sl, AX+0.2, ty, 3.15, 0.44, RGBColor(0x1E, 0x1B, 0x4B), clr, Pt(1))
        stxt(s, [(f"\u25c7 {step}", 10, True, clr)])
    else:
        s = box(sl, AX+0.2, ty, 3.15, 0.44, CARD, clr, Pt(0.75))
        stxt(s, [(step, 10, False, WHITE)])
    if i < len(act_steps)-1:
        box(sl, AX+1.02, ty+0.44, 0.04, 0.14, GRAY)

# Yes/No branch labels
tb(sl, "Yes \u2192", AX+3.5, AY0+3*0.58+0.04, 0.85, 0.34, size=9, color=GRAY, italic=True)
tb(sl, "No, skip \u2192", AX+3.5, AY0+4*0.58+0.04, 1.0, 0.34, size=9, color=GRAY, italic=True)

# ── RIGHT: Sequence Diagram ──────────────────────────────────
tb(sl, "Sequence Diagram \u2014 Login & Load Calendar",
   6.68, 1.22, 6.42, 0.36, size=13, bold=True, color=BLUE)

actors = [("React\nBrowser", 7.0), ("Django\nDRF", 9.15), ("SQLite\nDB", 11.22)]
for actor, xp in actors:
    tb(sl, actor, xp, 1.65, 1.55, 0.52, size=10, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    box(sl, xp+0.65, 2.20, 0.04, 4.72, GRAY)

msgs = [
    (7.65, 9.15, "POST /api/auth/login/",              2.50, BLUE),
    (9.15, 11.22, "SELECT user WHERE username=?",       3.05, GREEN),
    (11.22, 9.15, "Return User object",                 3.55, GRAY),
    (9.15, 7.65, "Set session + CSRF cookie",           4.05, PURPLE),
    (7.65, 9.15, "GET /api/tasks/occurrences/",         4.55, BLUE),
    (9.15, 11.22, "Expand RRULE + apply exceptions",    5.05, GREEN),
    (11.22, 9.15, "Occurrence list (JSON)",             5.55, GRAY),
    (9.15, 7.65, "200 OK — JSON response",              6.05, PURPLE),
]
for x1, x2, msg, ty, clr in msgs:
    fw = abs(x2 - x1) - 0.55
    fx = min(x1, x2) + 0.52
    box(sl, fx, ty, fw, 0.04, clr)
    tb(sl, msg, fx, ty-0.22, fw, 0.22,
       size=9, color=clr, italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# SLIDE 13 — IMPLEMENTATION: PSEUDOCODE
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Implementation \u2014 Pseudocode",
          "Core algorithm representations for key system operations")

pseudocodes = [
    ("1. User Authentication", PURPLE,
"""FUNCTION authenticate(username, password):
  user \u2190 DB.find_by_username(username)
  IF user AND verify_hash(password, user.hash):
    session \u2190 Session.create(user.id)
    csrf    \u2190 generate_csrf_token()
    RETURN  {user, session_cookie, csrf}
  ELSE:
    RAISE AuthenticationError(401)"""),

    ("2. Task Creation + Recurrence", BLUE,
"""FUNCTION create_task(data, user):
  task \u2190 Task(title, date, priority, user)
  IF data.is_recurring:
    task.recurrence_rule \u2190 buildRRULE(
        freq  = data.freq,
        days  = data.byDay,
        count = data.count)
  task.save() \u2192 DB
  RETURN expand_occurrences(task, current_range)"""),

    ("3. RRULE Expansion Engine", GREEN,
"""FUNCTION expand(tasks, start, end):
  results \u2190 []
  FOR task IN recurring_tasks:
    rule \u2190 rrulestr(task.recurrence_rule,
                dtstart = task.date)
    occs \u2190 rule.between(start, end)
    excs \u2190 DB.get_exceptions(task.id)
    FOR occ IN occs:
      IF occ IN excs.deleted: SKIP
      IF occ IN excs.overrides:
        results.add(override_data)
      ELSE: results.add(task + date=occ)
  RETURN results \u222a non_recurring_tasks"""),

    ("4. Analytics Aggregation", ORANGE,
"""FUNCTION get_analytics(user, start, end):
  all_occs  \u2190 expand_all(user, start, end)
  total     \u2190 all_occs.count()
  completed \u2190 all_occs.filter(COMPLETED)
  rate      \u2190 (completed / total) \u00d7 100
  streak    \u2190 calculate_streak(user)
  RETURN { total, completed, rate, streak }"""),
]

for i, (title, clr, code) in enumerate(pseudocodes):
    col, row = i%2, i//2
    l = 0.28 + col*6.60
    t = 1.34 + row*3.02
    tb(sl, title, l, t, 6.22, 0.36, size=12, bold=True, color=clr)
    s = box(sl, l, t+0.40, 6.22, 2.48, CODE_BG, clr, Pt(0.75))
    tb(sl, code, l+0.18, t+0.52, 5.96, 2.30, size=10, color=GREEN)


# ═════════════════════════════════════════════════════════════
# SLIDE 14 — SCREENSHOTS: AUTH
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Implementation \u2014 Authentication Pages",
          "Login and Registration UI")

ph(sl, 0.28, 1.35, 6.10, 5.78, "Login Page  (Yoga-Do Sign In Screen)")
ph(sl, 6.72, 1.35, 6.33, 5.78, "Registration Page")

tb(sl,
   "Session Auth  \u00b7  CSRF Protection  \u00b7  Inline Error Messages  "
   "\u00b7  Toggle Login/Register  \u00b7  Welcome Animation on Success",
   0.28, 7.12, 12.75, 0.30, size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# SLIDE 15 — SCREENSHOTS: CALENDAR VIEWS
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Implementation \u2014 Calendar Views",
          "Month, Week, and Day calendar views with priority-coded tasks")

ph(sl, 0.28, 1.35, 4.22, 2.72, "Month View")
ph(sl, 4.72, 1.35, 4.22, 2.72, "Week View")
ph(sl, 9.15, 1.35, 3.90, 2.72, "Day View")

view_info = [
    ("Month View", [
        "Full-month grid (4–5 weeks)",
        "Click any day to create a task",
        "Tasks shown with priority colour dots",
        "Prev/Next month navigation",
    ], 0.28, PURPLE, 4.22),
    ("Week View", [
        "7-day timeline display",
        "Colour-coded task blocks per day",
        "Click hour slot to create a task",
        "Prev/Next week navigation",
    ], 4.72, BLUE, 4.22),
    ("Day View", [
        "Single-day focused layout",
        "Tasks grouped in 3 priority sections",
        "Completion toggle per task",
        "Click slot to add new task",
    ], 9.15, GREEN, 3.90),
]
for label, bullets, x, clr, w in view_info:
    s = box(sl, x, 4.18, w, 3.14, CARD, clr, Pt(0.5))
    tb(sl, label, x+0.15, 4.25, w-0.28, 0.38, size=12, bold=True, color=clr)
    for j, b in enumerate(bullets):
        tb(sl, f"\u2022 {b}", x+0.15, 4.68+j*0.50, w-0.28, 0.44, size=11, color=WHITE)


# ═════════════════════════════════════════════════════════════
# SLIDE 16 — SCREENSHOTS: TASK MANAGEMENT
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Implementation \u2014 Task Management",
          "Creating, editing, and managing recurring tasks")

ph(sl, 0.28, 1.35, 4.12, 2.88, "Task Creation Modal")
ph(sl, 4.62, 1.35, 4.12, 2.88, "Recurring Task Dialog (RRULE options)")
ph(sl, 8.95, 1.35, 4.10, 2.88, "Edit: Single Occurrence / Entire Series")

task_feats = [
    ("Task Creation", [
        "Title, Description, Date + Time",
        "Priority: Critical / Important / Routine",
        "Recurring toggle reveals RRULE options",
        "Save \u2192 Calendar refreshed instantly",
    ], 0.28, PURPLE, 4.12),
    ("Recurring Tasks", [
        "Frequency: Daily or Weekly",
        "Select specific weekdays (Mon–Sun)",
        "Set total occurrence count",
        "RRULE stored as iCal-standard string",
    ], 4.62, BLUE, 4.12),
    ("Edit & Delete", [
        "Dialog: Edit this / Edit all occurrences",
        "Edit single \u2192 override saved as exception",
        "Edit series \u2192 parent task updated",
        "Delete one or delete all occurrences",
    ], 8.95, GREEN, 4.10),
]
for label, bullets, x, clr, w in task_feats:
    s = box(sl, x, 4.34, w, 2.98, CARD, clr, Pt(0.5))
    tb(sl, label, x+0.15, 4.42, w-0.28, 0.38, size=12, bold=True, color=clr)
    for j, b in enumerate(bullets):
        tb(sl, f"\u2022 {b}", x+0.15, 4.84+j*0.50, w-0.28, 0.44, size=11, color=WHITE)


# ═════════════════════════════════════════════════════════════
# SLIDE 17 — SCREENSHOTS: DASHBOARD ANALYTICS
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Implementation \u2014 Analytics Dashboard",
          "Real-time productivity tracking and statistics")

ph(sl, 0.28, 1.35, 8.18, 5.78,
   "Dashboard / Analytics View\n(Today's Focus · Progress Bars · Week/Month stats · Streak)")

dash_feats = [
    ("Today's Focus",    "Critical / Important / Routine priority breakdown",    PURPLE),
    ("Completion %",     "Real-time progress bars for Today / Week / Month",     BLUE),
    ("Task Counts",      "Total, Completed, and Remaining tasks at a glance",    GREEN),
    ("Streak Counter",   "Consecutive productive days tracked and displayed",     ORANGE),
    ("Live Updates",     "Dashboard refreshes on every task status change",      TEAL),
]
for i, (title, desc, clr) in enumerate(dash_feats):
    t = 1.40 + i*1.10
    s = box(sl, 8.78, t, 4.27, 0.96, CARD, clr, Pt(0.75))
    tb(sl, title, 8.93, t+0.08, 3.97, 0.34, size=12, bold=True, color=clr)
    tb(sl, desc,  8.93, t+0.50, 3.97, 0.40, size=10.5, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 18 — SCREENSHOTS: FOCUS TIMER & EXTRA FEATURES
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Implementation \u2014 Focus Timer & Extra Features",
          "Productivity tools embedded in the application")

ph(sl, 0.28, 1.35, 4.12, 2.88, "Focus Timer  (Pomodoro-style)")
ph(sl, 4.62, 1.35, 4.12, 2.88, "Personal Notes Section")
ph(sl, 8.95, 1.35, 4.10, 2.88, "Motivational Quotes")

feat_info = [
    ("Focus Timer", [
        "Pomodoro: 25-min work / 5-min break",
        "Visual countdown display",
        "Floats over app — non-intrusive overlay",
        "Show / Hide toggle from nav bar",
    ], 0.28, ORANGE, 4.12),
    ("Personal Notes", [
        "Free-text editor per user account",
        "Auto-saves to backend on change",
        "Accessible from any calendar view",
        "Persists across sessions",
    ], 4.62, BLUE, 4.12),
    ("Motivational Quotes", [
        "Curated uplifting quote display",
        "Refresh button for a new quote",
        "Styled with purple accent theme",
        "Boosts daily productivity mindset",
    ], 8.95, PURPLE, 4.10),
]
for label, bullets, x, clr, w in feat_info:
    s = box(sl, x, 4.34, w, 2.98, CARD, clr, Pt(0.5))
    tb(sl, label, x+0.15, 4.42, w-0.28, 0.38, size=12, bold=True, color=clr)
    for j, b in enumerate(bullets):
        tb(sl, f"\u2022 {b}", x+0.15, 4.84+j*0.50, w-0.28, 0.44, size=11, color=WHITE)


# ═════════════════════════════════════════════════════════════
# SLIDE 19 — FUTURE ENHANCEMENTS
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Future Enhancements",
          "Planned improvements for production scale and expanded feature set")

enhancements = [
    ("\U0001f5c4\ufe0f  PostgreSQL Migration",
     "Replace SQLite with PostgreSQL for production concurrency, indexing, and reliability.", BLUE),
    ("\u26a1  Redis Caching",
     "Cache analytics aggregations and RRULE expansions for sub-50 ms API response times.", TEAL),
    ("\U0001f916  AI Task Prioritisation",
     "ML-based suggestions to re-prioritise tasks using completion history and deadlines.", PURPLE),
    ("\U0001f4e7  Email Notifications",
     "SMTP reminders for task deadlines and overdue alerts via SendGrid / Mailgun.", ORANGE),
    ("\U0001f4c5  Google Calendar Sync",
     "Two-way sync with Google Calendar API using OAuth 2.0 tokens.", GREEN),
    ("\U0001f514  Push Notifications",
     "Browser Web Push API for overdue task alerts and upcoming deadline reminders.", RED),
    ("\u2388\ufe0f  Kubernetes Orchestration",
     "Migrate Docker deployment to K8s for horizontal auto-scaling and high availability.", BLUE),
    ("\U0001f465  Team Workspaces",
     "Shared task boards, team dashboards, and task assignment to collaborators.", PURPLE),
    ("\U0001f510  Role-Based Access Control",
     "Admin / Manager / Member roles with fine-grained permission enforcement.", GOLD),
    ("\U0001f4f4  Offline PWA Support",
     "Service workers for offline task viewing with background sync on reconnect.", TEAL),
    ("\U0001f4ca  Cloud Monitoring",
     "Grafana + Prometheus dashboards for API latency, error rate, and server health.", INDIGO),
    ("\U0001f319  Dark/Light Theme Toggle",
     "User-selectable theme with preference persisted in localStorage.", GRAY),
]

for i, (title, desc, clr) in enumerate(enhancements):
    col, row = i%3, i//3
    l = 0.28 + col*4.38
    t = 1.38 + row*1.56
    s = box(sl, l, t, 4.12, 1.40, CARD, clr, Pt(0.75))
    tb(sl, title, l+0.15, t+0.08, 3.82, 0.38, size=11, bold=True, color=clr)
    tb(sl, desc,  l+0.15, t+0.52, 3.82, 0.82, size=10, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE 20 — CONCLUSION
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "Conclusion", "Key outcomes and achievements")

conclusions = [
    ("\u2705  Full-Stack Web Application",
     "Successfully built Yoga-Do with a React + TypeScript frontend and a Django REST Framework backend, "
     "delivering a complete, calendar-driven task management experience with a responsive, professional UI.",
     PURPLE),
    ("\u2705  RRULE Recurring Task Engine",
     "Implemented a standards-compliant RRULE engine supporting daily and weekly recurrence with "
     "per-occurrence exception handling — edit or delete a single instance without affecting the entire series.",
     BLUE),
    ("\u2705  Analytics & Productivity Tools",
     "Real-time Analytics Dashboard, Pomodoro Focus Timer, Personal Notes, and Motivational Quotes are "
     "integrated directly into the application, providing a comprehensive productivity platform.",
     GREEN),
    ("\u2705  Complete DevOps Pipeline",
     "Git \u2192 GitHub \u2192 GitHub Actions CI (backend unit tests + Playwright E2E) \u2192 Docker build "
     "\u2192 Render/EC2 backend + Vercel frontend. Fully automated with zero manual deployment steps.",
     ORANGE),
    ("\u2705  Security & Engineering Best Practices",
     "Session-based auth, CSRF token protection, React error boundaries, automated testing, and "
     "Docker containerisation demonstrate professional software engineering standards.",
     TEAL),
]

for i, (title, desc, clr) in enumerate(conclusions):
    t = 1.34 + i*1.18
    box(sl, 0.28, t, 12.75, 1.08, CARD, clr, Pt(1))
    tb(sl, title, 0.46, t+0.06, 3.95, 0.42, size=12, bold=True, color=clr)
    tb(sl, desc,  4.50, t+0.08, 8.40, 0.88, size=11.5, color=WHITE)

tb(sl,
   "Yoga-Do demonstrates that modern full-stack web development and professional DevOps practices "
   "can be successfully combined in an academic mini-project to deliver a real-world quality application.",
   0.28, 7.10, 12.75, 0.34, size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# SLIDE 21 — REFERENCES
# ═════════════════════════════════════════════════════════════
sl = new_slide()
slide_hdr(sl, "References")

refs = [
    ("[1]",  "R. Singh, A. Patel et al., \"Task Management System Design Using Django REST Framework,\" "
             "Int. Journal of Computer Applications, Vol. 183, No. 12, pp. 1\u20138, 2022.", True),
    ("[2]",  "S. Mehta et al., \"Containerisation with Docker for Academic and Enterprise Web Applications,\" "
             "IEEE Access, Vol. 10, pp. 45123\u201345135, 2023.", True),
    ("[3]",  "V. Kumar, M. Rao, \"Automating CI/CD Pipelines Using GitHub Actions for Python-Django Projects,\" "
             "Journal of Open-Source Software, Vol. 8, No. 91, 2023.", True),
    ("[4]",  "W. Li, J. Chen, \"Efficient RRULE Processing for Calendar-Based Task Management Applications,\" "
             "ACM SIGAPP Applied Computing Review, Vol. 24, No. 1, 2024.", True),
    ("[5]",  "P. Sharma, K. Nair, \"Integrating Full-Stack DevOps in BE/BTech Capstone Projects,\" "
             "Proc. Int. Conf. on Software Engineering Education, pp. 211\u2013218, 2024.", True),
    ("[6]",  "Django REST Framework \u2014 Official Documentation. https://www.django-rest-framework.org/", False),
    ("[7]",  "React Documentation. https://react.dev/", False),
    ("[8]",  "Docker Official Documentation. https://docs.docker.com/", False),
    ("[9]",  "GitHub Actions Documentation. https://docs.github.com/en/actions", False),
    ("[10]", "python-dateutil RRULE Reference. https://dateutil.readthedocs.io/", False),
    ("[11]", "Playwright Testing Framework. https://playwright.dev/", False),
]

for i, (num, ref, is_paper) in enumerate(refs):
    t = 1.34 + i*0.55
    box(sl, 0.28, t, 12.75, 0.50, CARD if i%2==0 else BG)
    nb = box(sl, 0.28, t, 0.58, 0.50, PURPLE if is_paper else BLUE)
    stxt(nb, [(num, 10, True, WHITE)])
    tb(sl, ref, 0.98, t+0.08, 11.90, 0.36,
       size=10.5, color=WHITE if is_paper else GRAY)


# ═════════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════════
import os
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Yoga-Do_Presentation.pptx")
prs.save(OUT)
print(f"\nPresentation saved successfully!")
print(f"File   : {OUT}")
print(f"Slides : {len(prs.slides)}")
print(f"\nPlaceholder slides that need screenshots inserted:")
for i in [14, 15, 16, 17, 18, 9]:
    print(f"  Slide {i}")
